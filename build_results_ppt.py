#!/usr/bin/env python3
"""Build the IEEE WCCI 2026 MicroRTS LLM Competition results deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = "docs/MicroRTS_LLM_Competition_Results.pptx"

NAVY = RGBColor(0x10, 0x2A, 0x54)
ACCENT = RGBColor(0xE3, 0x6B, 0x18)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
DARK = RGBColor(0x1F, 0x29, 0x3A)
MUTED = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
SILVER = RGBColor(0x9A, 0xA0, 0xA6)
BRONZE = RGBColor(0xB0, 0x6A, 0x2C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.9), NAVY)
    add_rect(slide, 0, Inches(0.9), prs.slide_width, Inches(0.06), ACCENT)
    add_text(slide, Inches(0.4), Inches(0.15), Inches(11), Inches(0.6),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(11), Inches(0.35),
                 subtitle, size=12, color=RGBColor(0xCF, 0xD8, 0xE3))


def add_footer(slide, page):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             "2026 IEEE WCCI · MicroRTS LLM Game AI Competition",
             size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{page}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ===== Slide 1: Title =====
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
add_rect(s, 0, Inches(4.0), prs.slide_width, Inches(0.08), ACCENT)
add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.2),
         "MicroRTS LLM Game AI Competition", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.8),
         "Final Results & Analysis", size=28, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
         "2026 IEEE World Congress on Computational Intelligence",
         size=20, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(4.9), Inches(12), Inches(0.4),
         "Organizers: Chang Liu et al.   ·   13 teams   ·   15 agents evaluated",
         size=16, color=WHITE)
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.3),
         "github.com/drchangliu/MicroRTS  ·  drchangliu.github.io/MicroRTS",
         size=12, color=RGBColor(0xCF, 0xD8, 0xE3))

# ===== Slide 2: Background & History =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Background & History",
           "MicroRTS — a decade-old benchmark for real-time strategy AI")

# Left column: about MicroRTS
add_text(s, Inches(0.5), Inches(1.1), Inches(6.3), Inches(0.4),
         "About MicroRTS", size=16, bold=True, color=NAVY)
about_lines = [
    "• Lightweight Java RTS designed for AI research — workers, light / heavy / ranged "
    "units, bases, barracks, resources.",
    "• Created by Santiago Ontañón (Drexel) and presented at AIIDE 2013 (Combinatorial "
    "Multi-Armed Bandit work).",
    "• Fully observable or partially observable; deterministic engine; fast simulation "
    "(thousands of games per minute) — ideal for tree search, RL, and benchmark studies.",
    "• Maintained jointly with Levi Lelis and Rubens O. Moraes, who have driven the "
    "competition for years.",
    "• Open source: github.com/santiontanon/microrts",
]
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.2), Inches(3.0))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(about_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run(); r.text = l; r.font.size = Pt(12); r.font.color.rgb = DARK

# Right column: 32x32 screenshot
add_text(s, Inches(7.2), Inches(1.1), Inches(5.6), Inches(0.4),
         "Same engine, ten years on", size=16, bold=True, color=NAVY)
s.shapes.add_picture("docs/screenshots/32x32.png",
                     Inches(7.2), Inches(1.5), Inches(5.6), Inches(4.0))
add_text(s, Inches(7.2), Inches(5.5), Inches(5.6), Inches(0.3),
         "32×32 mid-play: LightRush vs HeavyRush at tick 800",
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom row: competition history timeline
add_text(s, Inches(0.5), Inches(4.7), Inches(6.3), Inches(0.4),
         "Competition history", size=16, bold=True, color=NAVY)
history = [
    ("2017", "IEEE-CIG", "First MicroRTS AI Competition"),
    ("2018", "IEEE-CIG", "Tree-search + scripted hybrids dominate"),
    ("2019", "IEEE-COG", "Renamed conference; CoacAI emerges"),
    ("2020-21", "IEEE-COG", "POLightRush, Mayari, Tiamat refine the ladder"),
    ("2022-23", "IEEE-COG", "RL-track and PO-track expand"),
    ("2026", "IEEE-WCCI", "First LLM-only edition (this competition)"),
]
y0 = Inches(5.15)
for i, (year, venue, note) in enumerate(history):
    y = y0 + Inches(i * 0.28)
    add_rect(s, Inches(0.5), y, Inches(0.9), Inches(0.26),
             ACCENT if i == 5 else NAVY)
    add_text(s, Inches(0.55), y + Inches(0.02), Inches(0.85), Inches(0.22),
             year, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), y + Inches(0.02), Inches(1.2), Inches(0.22),
             venue, size=10, bold=True, color=DARK)
    add_text(s, Inches(2.7), y + Inches(0.02), Inches(4.2), Inches(0.22),
             note, size=10, color=DARK)
add_footer(s, 2)

# ===== Slide 3: Why an LLM Edition? =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Why an LLM Edition?",
           "Motivation for the 2026 WCCI track")

add_text(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6),
         "Use LLMs as the core of intelligence to solve RTS-game problems.",
         size=22, bold=True, color=NAVY)

add_text(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(0.9),
         "LLM-based RTS agents are not the top solution today — but the area has "
         "high growth potential, and is likely a research direction that compounds "
         "quickly over the next few years.",
         size=15, color=DARK)

# Two prongs side by side
add_text(s, Inches(0.5), Inches(2.85), Inches(12), Inches(0.4),
         "A two-pronged research program", size=16, bold=True, color=ACCENT)

# Prong 1
add_rect(s, Inches(0.5), Inches(3.35), Inches(6.1), Inches(3.4), LIGHT)
add_rect(s, Inches(0.5), Inches(3.35), Inches(6.1), Inches(0.7), NAVY)
add_text(s, Inches(0.7), Inches(3.45), Inches(5.7), Inches(0.55),
         "Prong 1 · How best to use an LLM to solve RTS?",
         size=15, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(4.15), Inches(5.7), Inches(0.4),
         "Architecture / scaffolding question",
         size=12, bold=True, color=ACCENT)
prong1 = [
    "• Where does the LLM sit in the loop — every tick, every N ticks, or only at strategy switches?",
    "• How do you compose an LLM with classical search (MCTS) or scripted policies?",
    "• What prompt format and game-state encoding makes the LLM reliable?",
    "• How do you handle latency, validity constraints, and JSON parsing failures?",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(5.7), Inches(2.1))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(prong1):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run(); r.text = l; r.font.size = Pt(11); r.font.color.rgb = DARK

# Prong 2
add_rect(s, Inches(6.75), Inches(3.35), Inches(6.1), Inches(3.4), LIGHT)
add_rect(s, Inches(6.75), Inches(3.35), Inches(6.1), Inches(0.7), NAVY)
add_text(s, Inches(6.95), Inches(3.45), Inches(5.7), Inches(0.55),
         "Prong 2 · Which LLM is best at RTS?",
         size=15, bold=True, color=WHITE)
add_text(s, Inches(6.95), Inches(4.15), Inches(5.7), Inches(0.4),
         "Benchmark / model-comparison question",
         size=12, bold=True, color=ACCENT)
prong2 = [
    "• A standardized RTS benchmark across open and closed models.",
    "• Measures planning, instruction-following, and structured-output reliability.",
    "• Holds scaffolding constant; isolates the model as the independent variable.",
    "• Tracks year-over-year LLM progress on a real-time, partially observable task.",
]
tb = s.shapes.add_textbox(Inches(6.95), Inches(4.55), Inches(5.7), Inches(2.1))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(prong2):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run(); r.text = l; r.font.size = Pt(11); r.font.color.rgb = DARK

add_footer(s, 3)

# ===== Slide 4: Competition Overview =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Competition Overview", "What we asked the teams to build")
add_text(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
         "Build an LLM-powered MicroRTS agent — no training, no fine-tuning, prompt-only.",
         size=18, bold=True, color=NAVY)

# Left: bullets
bullets = [
    "MicroRTS: open Java RTS benchmark (Ontañón et al.) — workers, light/heavy/ranged units, bases, barracks.",
    "Agents read the game state, prompt an LLM, parse JSON moves, submit them per tick.",
    "Success requires prompt engineering and strategy design, not weight updates.",
    "Submitted via GitHub PR; central server runs each agent against a six-anchor ladder.",
    "Tests LLM reasoning, planning, and instruction-following under real-time constraints.",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(7.5), Inches(4.4))
tf = tb.text_frame
tf.word_wrap = True
for i, b in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run(); r.text = "•  " + b
    r.font.size = Pt(14); r.font.color.rgb = DARK

# Right: 16x16 screenshot showcase
add_text(s, Inches(8.3), Inches(1.8), Inches(4.6), Inches(0.35),
         "MicroRTS · 16×16 · mid-play", size=12, bold=True, color=NAVY)
s.shapes.add_picture("docs/screenshots/16x16.png",
                     Inches(8.3), Inches(2.15), Inches(4.6), Inches(4.6))
add_text(s, Inches(8.3), Inches(6.75), Inches(4.6), Inches(0.3),
         "Blue = Player 0 · Pink = Player 1 · Orange circle = Light unit",
         size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Stats strip across the bottom
add_rect(s, Inches(0.5), Inches(6.3), Inches(7.5), Inches(0.5), ACCENT)
add_text(s, Inches(0.65), Inches(6.35), Inches(7.3), Inches(0.4),
         "13 teams  ·  15 binaries  ·  4 map sizes (8×8 → 64×64)  ·  6 anchors  ·  1 A+",
         size=13, bold=True, color=WHITE)
add_footer(s, 4)

# ===== Slide 3: Format & Scoring =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Tournament Format & Scoring", "Single-elimination, multi-map ladder")

# Maps table
add_text(s, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
         "Map gauntlet (per agent)", size=16, bold=True, color=NAVY)
maps = [("Map", "Cycle cap"),
        ("8×8  basesWorkers", "1,500"),
        ("16×16 basesWorkers", "3,000"),
        ("32×32 basesWorkers", "5,000"),
        ("64×64 GardenOfWar", "8,000")]
for i, (a, b) in enumerate(maps):
    y = Inches(1.55 + i * 0.42)
    bg = NAVY if i == 0 else (LIGHT if i % 2 else WHITE)
    fg = WHITE if i == 0 else DARK
    add_rect(s, Inches(0.5), y, Inches(4.0), Inches(0.42), bg,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(s, Inches(4.5), y, Inches(1.8), Inches(0.42), bg,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s, Inches(0.6), y + Inches(0.06), Inches(3.9), Inches(0.3), a,
             size=13, bold=(i == 0), color=fg)
    add_text(s, Inches(4.6), y + Inches(0.06), Inches(1.7), Inches(0.3), b,
             size=13, bold=(i == 0), color=fg, align=PP_ALIGN.RIGHT)

# Anchors table
add_text(s, Inches(6.8), Inches(1.1), Inches(6), Inches(0.4),
         "Anchor opponents (weight)", size=16, bold=True, color=NAVY)
anchors = [("Opponent", "Tier", "Weight"),
           ("RandomBiasedAI", "easy", "10"),
           ("WorkerRush", "medium", "15"),
           ("LightRush", "medium", "15"),
           ("HeavyRush", "medium-hard", "20"),
           ("Tiamat", "hard", "20"),
           ("CoacAI", "hard", "20")]
for i, row in enumerate(anchors):
    y = Inches(1.55 + i * 0.42)
    bg = NAVY if i == 0 else (LIGHT if i % 2 else WHITE)
    fg = WHITE if i == 0 else DARK
    add_rect(s, Inches(6.8), y, Inches(3.0), Inches(0.42), bg,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(s, Inches(9.8), y, Inches(1.7), Inches(0.42), bg,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(s, Inches(11.5), y, Inches(1.3), Inches(0.42), bg,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s, Inches(6.9), y + Inches(0.06), Inches(2.9), Inches(0.3),
             row[0], size=13, bold=(i == 0), color=fg)
    add_text(s, Inches(9.9), y + Inches(0.06), Inches(1.6), Inches(0.3),
             row[1], size=13, bold=(i == 0), color=fg)
    add_text(s, Inches(11.6), y + Inches(0.06), Inches(1.1), Inches(0.3),
             row[2], size=13, bold=(i == 0), color=fg, align=PP_ALIGN.RIGHT)

# Scoring rule
add_rect(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.9), LIGHT)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
         "Scoring", size=16, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.4),
         "Weighted points per win (1.0) and draw (0.5), averaged across maps. "
         "A loss to an anchor stops that map (single-elimination) and the "
         "agent banks whatever it earned reaching that anchor.\n\n"
         "Grades:  A+ ≥ 90    ·    A ≥ 80    ·    B ≥ 70    ·    C ≥ 60    ·    D ≥ 40    ·    F < 40\n"
         "Default server model: llama3.1:8b (Ollama). Teams could declare a preferred model.",
         size=14, color=DARK)
add_footer(s, 5)

# ===== Slide 4 (new): MicroRTS gameplay gallery =====
s = prs.slides.add_slide(BLANK)
add_header(s, "MicroRTS — Four Map Sizes",
           "Same engine, same rules, very different strategic horizon")

# 2x2 grid of screenshots
gallery = [
    ("docs/screenshots/8x8.png",   "8×8  basesWorkers",
     "1,500 cycles · WorkerRush mirror @ tick 180"),
    ("docs/screenshots/16x16.png", "16×16  basesWorkers",
     "3,000 cycles · LightRush mirror @ tick 600 (HP bars = engaged)"),
    ("docs/screenshots/32x32.png", "32×32  basesWorkers",
     "5,000 cycles · LightRush vs HeavyRush @ tick 800"),
    ("docs/screenshots/64x64.png", "64×64  GardenOfWar",
     "8,000 cycles · LightRush vs RangedRush @ tick 1500 — walls split the map"),
]
positions = [(0.5, 1.1), (6.95, 1.1), (0.5, 4.25), (6.95, 4.25)]
for (x, y), (path, label, sub) in zip(positions, gallery):
    add_rect(s, Inches(x), Inches(y), Inches(5.85), Inches(2.6), WHITE,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    s.shapes.add_picture(path, Inches(x + 0.05), Inches(y + 0.05),
                         Inches(2.55), Inches(2.5))
    add_text(s, Inches(x + 2.75), Inches(y + 0.15), Inches(3.0), Inches(0.4),
             label, size=14, bold=True, color=NAVY)
    add_text(s, Inches(x + 2.75), Inches(y + 0.6), Inches(3.0), Inches(1.8),
             sub, size=10, color=DARK)
add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         "Screenshots captured mid-play with scripted rush bots on the organizers' server.",
         size=10, color=MUTED)
add_footer(s, 6)

# ===== Slide 5: Final Leaderboard =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Final Leaderboard", "Best score per team across all tournament runs")

leaderboard = [
    (1, "AlliBot",        100.5, "A+", "qwen3:14b",   "Hybrid + LLM advisor/search"),
    (2, "Mayari-LLM",      87.0, "A",  "llama3.1:8b", "MayariBot + LLM consultation"),
    (3, "HOPE",            79.0, "B",  "llama3.1:8b", "Ollama + MCTS + prediction"),
    (4, "Chase",           69.0, "C",  "llama3.1:8b", "Rush engine + advisor every ~500 ticks"),
    (5, "yebot",           69.0, "C",  "qwen3:8b",    "Macro-LLM + hard-coded micro"),
    (6, "Fortress Bot",    57.5, "D",  "llama3.1:8b", "LLM planner + scripted opening + rush mirrors"),
    (7, "PenguinBot",      54.0, "D",  "llama3.1:8b", "NaiveMCTS + LLM stance controller"),
    (8, "xiebot",          54.0, "D",  "llama3.1:8b", "LLM planner + selective tactical MCTS"),
    (9, "jmurr",           35.0, "F",  "none",        "Symmetric self/opponent state eval"),
    (10, "AI4PC",          33.8, "F",  "llama3.1:8b", "LLM picks rush every 100 ticks (counter-triangle)"),
    (11, "Parker's Bot",   33.0, "F",  "llama3.1:8b", "LLM-guided adaptive rush + fallback"),
    (12, "jmurrllm",       24.0, "F",  "llama3.1:8b", "Light-rush opener + LLM classifier + counter table"),
    (13, "EAGLE",           3.8, "F",  "llama3.1:8b", "Structured prompt policy, role-based workers"),
    (13, "Adil Bot",        3.8, "F",  "none",        "Rule-based: harvest/build/train + early rush"),
]

# Header row
hdr_y = Inches(1.1)
cols = [(0.4, 0.6, "#"),
        (1.0, 2.4, "Team"),
        (3.4, 1.1, "Score"),
        (4.5, 0.8, "Grade"),
        (5.3, 1.9, "Model"),
        (7.2, 5.9, "Approach")]
for x, w, label in cols:
    add_rect(s, Inches(x), hdr_y, Inches(w), Inches(0.38), NAVY)
    add_text(s, Inches(x + 0.08), hdr_y + Inches(0.06), Inches(w - 0.1),
             Inches(0.3), label, size=12, bold=True, color=WHITE)

row_h = 0.36
for i, (rk, name, sc, gr, mdl, appr) in enumerate(leaderboard):
    y = Inches(1.5 + i * row_h)
    bg = LIGHT if i % 2 == 0 else WHITE
    medal = None
    if rk == 1:
        bg, medal = RGBColor(0xFD, 0xF3, 0xD0), GOLD
    elif rk == 2:
        bg, medal = RGBColor(0xEE, 0xEE, 0xEE), SILVER
    elif rk == 3:
        bg, medal = RGBColor(0xF6, 0xE4, 0xD0), BRONZE
    for x, w, _ in cols:
        add_rect(s, Inches(x), y, Inches(w), Inches(row_h), bg,
                 line=RGBColor(0xE5, 0xE7, 0xEA))
    add_text(s, Inches(cols[0][0] + 0.08), y + Inches(0.05), Inches(cols[0][1] - 0.1),
             Inches(0.3), str(rk), size=11, bold=True,
             color=(medal or DARK))
    add_text(s, Inches(cols[1][0] + 0.08), y + Inches(0.05), Inches(cols[1][1] - 0.1),
             Inches(0.3), name, size=11, bold=(rk <= 3), color=DARK)
    add_text(s, Inches(cols[2][0] + 0.08), y + Inches(0.05), Inches(cols[2][1] - 0.1),
             Inches(0.3), f"{sc:.1f}", size=11, bold=True, color=DARK,
             align=PP_ALIGN.RIGHT)
    grade_color = {"A+": RGBColor(0x1E, 0x88, 0x4E),
                   "A":  RGBColor(0x1E, 0x88, 0x4E),
                   "B":  RGBColor(0x33, 0x6F, 0xB8),
                   "C":  RGBColor(0xE3, 0x6B, 0x18),
                   "D":  RGBColor(0xC4, 0x46, 0x21),
                   "F":  RGBColor(0x8E, 0x16, 0x16)}[gr]
    add_text(s, Inches(cols[3][0] + 0.08), y + Inches(0.05), Inches(cols[3][1] - 0.1),
             Inches(0.3), gr, size=11, bold=True, color=grade_color,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(cols[4][0] + 0.08), y + Inches(0.05), Inches(cols[4][1] - 0.1),
             Inches(0.3), mdl, size=10, color=DARK)
    add_text(s, Inches(cols[5][0] + 0.08), y + Inches(0.05), Inches(cols[5][1] - 0.1),
             Inches(0.3), appr, size=10, color=DARK)
add_footer(s, 7)

# ===== Slide 6: Winner Spotlight - AlliBot =====
s = prs.slides.add_slide(BLANK)
add_header(s, "🏆 Champion: AlliBot", "Score 100.5 · Grade A+ · Model qwen3:14b")
add_text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "The only A+ submission, and the only agent to complete the full 4-map gauntlet.",
         size=18, bold=True, color=NAVY)

# Architecture column
add_rect(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(4.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.95), Inches(5.8), Inches(0.4),
         "Architecture", size=16, bold=True, color=NAVY)
arch_lines = [
    "• Mayari-derived hybrid macro/micro execution layer",
    "• Rush-defense rule set: counter-build vs. visible enemy composition",
    "• Small-map (4×4 / 8×8) WorkerRush delegation",
    "• qwen3:14b advisor: high-level strategic guidance",
    "• Search component biased by LLM priors over actions",
    "• Deterministic fallback if LLM stalls (latency budget)",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(2.35), Inches(5.8), Inches(3.9))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(arch_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run(); r.text = l; r.font.size = Pt(14); r.font.color.rgb = DARK

# Per-map result
add_rect(s, Inches(6.8), Inches(1.85), Inches(6.0), Inches(4.5), LIGHT)
add_text(s, Inches(7.0), Inches(1.95), Inches(5.8), Inches(0.4),
         "Per-map performance", size=16, bold=True, color=NAVY)
map_rows = [("Map", "Score", "Grade", "Eliminated by"),
            ("8×8",   "69",  "C",  "Tiamat"),
            ("16×16", "117", "A+", "—"),
            ("32×32", "96",  "A+", "CoacAI"),
            ("64×64", "120", "A+", "—")]
for i, row in enumerate(map_rows):
    y = Inches(2.45 + i * 0.55)
    bg = NAVY if i == 0 else WHITE
    fg = WHITE if i == 0 else DARK
    widths = [1.2, 1.2, 1.2, 2.4]
    xs = [7.0, 8.2, 9.4, 10.6]
    for x, w in zip(xs, widths):
        add_rect(s, Inches(x), y, Inches(w), Inches(0.55), bg,
                 line=RGBColor(0xDD, 0xDD, 0xDD))
    for j, val in enumerate(row):
        bold = (i == 0) or (j == 2 and val.startswith("A"))
        color = fg
        if i > 0 and j == 2 and val == "A+":
            color = RGBColor(0x1E, 0x88, 0x4E)
        add_text(s, Inches(xs[j] + 0.1), y + Inches(0.13),
                 Inches(widths[j] - 0.2), Inches(0.3),
                 val, size=13, bold=bold, color=color)

add_rect(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.45), ACCENT)
add_text(s, Inches(0.65), Inches(6.5), Inches(12), Inches(0.35),
         "Headline: AlliBot is the only entry to beat CoacAI on any map "
         "(32×32, 2W–1L) and to score A+ on three of four maps.",
         size=13, bold=True, color=WHITE)
add_footer(s, 8)

# ===== Slide 6: Submission Archetypes =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Five Design Patterns", "How teams composed LLMs with classical game AI")

archetypes = [
    ("1.  MCTS + LLM stance",
     "Search-led, LLM-advised",
     "PenguinBot · HOPE · xiebot",
     "Tree search drives micro; LLM is consulted every few hundred ticks "
     "to bias exploration toward ATTACK/DEFEND or to set macro priors."),
    ("2.  LLM strategy controller over scripted rushes",
     "Script-led, LLM-routed",
     "AI4PC · Parker's Bot · Chase · Fortress Bot · jmurrllm",
     "Each tick is executed by a canonical rush bot (Worker/Light/Heavy/Ranged); "
     "LLM picks which rush — often with a counter-triangle fallback."),
    ("3.  Hybrid macro/micro with LLM macro layer",
     "LLM-led, scripted micro",
     "AlliBot · Mayari-LLM · yebot",
     "LLM proposes goals; deterministic code handles harvesting, building, "
     "and combat micro. Top performers cluster here."),
    ("4.  Pure-prompt LLM policy",
     "LLM-only",
     "EAGLE",
     "Every action generated from a structured prompt under rule constraints. "
     "Heavy latency cost — hard to clear even the weakest anchor."),
    ("5.  Heuristic / no-LLM baselines",
     "No LLM",
     "jmurr · Adil Bot",
     "Submitted to measure how much the LLM layer is actually worth."),
]
y = Inches(1.1)
for title, sub, teams, desc in archetypes:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.1), LIGHT,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.1), ACCENT)
    add_text(s, Inches(0.75), y + Inches(0.08), Inches(7), Inches(0.35),
             title, size=14, bold=True, color=NAVY)
    add_text(s, Inches(8.0), y + Inches(0.08), Inches(4.8), Inches(0.35),
             sub, size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(0.75), y + Inches(0.42), Inches(12), Inches(0.3),
             "Teams: " + teams, size=11, bold=True, color=DARK)
    add_text(s, Inches(0.75), y + Inches(0.7), Inches(12), Inches(0.4),
             desc, size=11, color=MUTED)
    y += Inches(1.2)
add_footer(s, 9)

# ===== Slide 7: What Worked =====
s = prs.slides.add_slide(BLANK)
add_header(s, "What Worked · What Didn't", "Lessons from 15 agents")

# What worked panel
add_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(5.6),
         RGBColor(0xE6, 0xF3, 0xEC))
add_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(0.5),
         RGBColor(0x1E, 0x88, 0x4E))
add_text(s, Inches(0.65), Inches(1.2), Inches(6), Inches(0.35),
         "✓  What worked", size=18, bold=True, color=WHITE)
worked = [
    ("Hybrid architectures dominated.",
     "Every A/B-grade entry kept the LLM out of the per-tick loop and used it as a strategic advisor (100–500 ticks)."),
    ("Stronger models translated to wins…",
     "…when the rest of the stack was sound. AlliBot (qwen3:14b) is the only A+. yebot (qwen3:8b) edges out llama3.1:8b peers."),
    ("Counter-triangle reasoning recurs.",
     "Light beats Ranged · Heavy beats Light · Ranged beats Heavy — appears in 5 entries, reliable at 8×8."),
    ("Deterministic fallback paths.",
     "Every top-half entry has a heuristic plan B if the LLM stalls or returns invalid JSON."),
]
y = Inches(1.75)
for h, t in worked:
    add_text(s, Inches(0.7), y, Inches(5.8), Inches(0.3), h,
             size=13, bold=True, color=DARK)
    add_text(s, Inches(0.7), y + Inches(0.32), Inches(5.8), Inches(0.7), t,
             size=11, color=DARK)
    y += Inches(1.2)

# What didn't panel
add_rect(s, Inches(6.7), Inches(1.1), Inches(6.1), Inches(5.6),
         RGBColor(0xFA, 0xE9, 0xE6))
add_rect(s, Inches(6.7), Inches(1.1), Inches(6.1), Inches(0.5),
         RGBColor(0x8E, 0x16, 0x16))
add_text(s, Inches(6.85), Inches(1.2), Inches(6), Inches(0.35),
         "✗  What didn't", size=18, bold=True, color=WHITE)
didnt = [
    ("Per-tick LLM prompting.",
     "EAGLE (pure-prompt policy) struggled with action validity and latency — could not reliably beat RandomBiasedAI."),
    ("Small-map tuning rarely transferred.",
     "Mayari-LLM scored A+ on 16×16 but D on 8×8. Map-specific heuristics broke under different cycle budgets."),
    ("Stickiness without re-evaluation.",
     "Per-100-tick strategy pickers (AI4PC) lost on larger maps where the opponent's composition shifted faster than polling."),
    ("Tiamat was the bottleneck.",
     "Tiamat eliminated 3 top-half entries. Only AlliBot beat it on any map."),
]
y = Inches(1.75)
for h, t in didnt:
    add_text(s, Inches(6.9), y, Inches(5.8), Inches(0.3), h,
             size=13, bold=True, color=DARK)
    add_text(s, Inches(6.9), y + Inches(0.32), Inches(5.8), Inches(0.7), t,
             size=11, color=DARK)
    y += Inches(1.2)
add_footer(s, 10)

# ===== Slide 8: Per-map breakdown for top finishers =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Per-Map Performance — Top Finishers",
           "Score per map · grade · anchor that eliminated the agent")

rows = [
    ("Team", "8×8", "16×16", "32×32", "64×64"),
    ("AlliBot",     "69 C\nTiamat",      "117 A+\n—",         "96 A+\nCoacAI",     "120 A+\n—"),
    ("Mayari-LLM",  "54 D\nWorkerRush",  "120 A+\n—",         "— not run",         "— not run"),
    ("HOPE",        "79 B\nadvanced",    "— not run",         "— not run",         "— not run"),
    ("Chase",       "69 C\nadvanced",    "— not run",         "— not run",         "— not run"),
    ("yebot",       "69 C\nadvanced",    "— not run",         "— not run",         "— not run"),
    ("Fortress Bot","79 B\nTiamat",      "36 F\nLightRush",   "— not run",         "— not run"),
    ("AI4PC",       "79 B\nTiamat",      "22 F\nHeavyRush",   "22 F\nHeavyRush",   "12 F\nHeavyRush"),
]
col_xs = [0.5, 3.2, 5.6, 8.0, 10.4]
col_w = [2.6, 2.3, 2.3, 2.3, 2.3]
row_h = 0.68
for i, row in enumerate(rows):
    y = Inches(1.2 + i * row_h)
    for j, val in enumerate(row):
        bg = NAVY if i == 0 else (LIGHT if i % 2 else WHITE)
        fg = WHITE if i == 0 else DARK
        add_rect(s, Inches(col_xs[j]), y, Inches(col_w[j]), Inches(row_h),
                 bg, line=RGBColor(0xDD, 0xDD, 0xDD))
        is_first = (j == 0)
        size = 12 if i == 0 else (12 if is_first else 11)
        bold = (i == 0) or is_first
        add_text(s, Inches(col_xs[j] + 0.1), y + Inches(0.06),
                 Inches(col_w[j] - 0.15), Inches(row_h - 0.1),
                 val, size=size, bold=bold, color=fg)

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         "Only AlliBot completed all four maps; only AlliBot beat CoacAI on any map.",
         size=11, color=MUTED)
add_footer(s, 11)

# ===== Slide 9: Headline Numbers =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Headline Numbers")

stats = [
    ("13", "teams submitted via GitHub PR"),
    ("15", "agent binaries evaluated"),
    ("1", "A+ submission (AlliBot, 100.5)"),
    ("38%", "of teams scored C or higher"),
    ("4", "different LLMs used: llama3.1:8b · qwen3:14b · qwen3:8b · gpt-5 (dev)"),
    ("3", "no-LLM/heuristic baselines (control group)"),
    ("1", "agent beat CoacAI on a map (AlliBot, 32×32)"),
]
# Layout as a 2x4 grid
cell_w, cell_h = Inches(6.1), Inches(1.4)
positions = [(0.5, 1.2), (6.7, 1.2),
             (0.5, 2.75), (6.7, 2.75),
             (0.5, 4.3), (6.7, 4.3),
             (0.5, 5.85)]
for (x, y), (num, label) in zip(positions, stats):
    add_rect(s, Inches(x), Inches(y), cell_w, cell_h, LIGHT,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(s, Inches(x), Inches(y), Inches(0.15), cell_h, ACCENT)
    add_text(s, Inches(x + 0.3), Inches(y + 0.15), Inches(2.5), Inches(1.1),
             num, size=44, bold=True, color=NAVY)
    add_text(s, Inches(x + 3.0), Inches(y + 0.45), Inches(3.0), Inches(0.7),
             label, size=14, color=DARK)
add_footer(s, 12)

# ===== Slide 10: Future Work / Acknowledgments =====
s = prs.slides.add_slide(BLANK)
add_header(s, "Future Work & Acknowledgments")

# Future work
add_text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
         "Toward the 2027 edition", size=18, bold=True, color=NAVY)
future = [
    "Multiple games per matchup (variance reduction) and per-anchor seeds.",
    "Standardize LLM latency budget reporting to compare hybrid designs fairly.",
    "Allow team-selected model on the server, paired with a fixed compute cap.",
    "Add stronger anchors (POLightRush, NaiveMCTS+) to better separate A+/A entries.",
    "Per-game replay viewer hosted on the leaderboard site for transparency.",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
for i, b in enumerate(future):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run(); r.text = "•  " + b
    r.font.size = Pt(15); r.font.color.rgb = DARK

# Acknowledgments
add_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), LIGHT)
add_text(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.4),
         "Acknowledgments", size=18, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(5.05), Inches(12), Inches(1.7),
         "• All 13 teams for their submissions and reproducible code.\n"
         "• Santi Ontañón, Levi Lelis, and Rubens O. Moraes for the open MicroRTS engine.\n"
         "• IEEE WCCI 2026 competition committee for hosting the track.\n"
         "• Submissions, logs, and per-game results: github.com/drchangliu/MicroRTS",
         size=14, color=DARK)
add_footer(s, 13)

# ===== Slide 11: Thank you =====
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
add_rect(s, 0, Inches(3.6), prs.slide_width, Inches(0.08), ACCENT)
add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(1.5),
         "Thank you", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.9), Inches(12), Inches(0.6),
         "Questions & discussion", size=24, color=RGBColor(0xCF, 0xD8, 0xE3),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
         "drchangliu.github.io/MicroRTS  ·  github.com/drchangliu/MicroRTS",
         size=16, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.1), Inches(12), Inches(0.4),
         "Submissions remain open — try the 2027 leaderboard at the same address.",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Wrote {OUT}  ({len(prs.slides)} slides)")
